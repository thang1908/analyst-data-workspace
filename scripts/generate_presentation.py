import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    COLOR_DARK_BG = RGBColor(15, 23, 42)      # #0F172A (Navy slate)
    COLOR_CARD_BG = RGBColor(241, 245, 249)   # #F1F5F9 (Light slate)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_MAIN = RGBColor(30, 41, 59)    # #1E293B
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    COLOR_PRIMARY_BLUE = RGBColor(37, 99, 235) # #2563EB
    COLOR_LIGHT_BLUE = RGBColor(239, 246, 255) # #EFF6FF
    COLOR_BORDER_BLUE = RGBColor(191, 219, 254) # #BFDBFE
    COLOR_SUCCESS_GREEN = RGBColor(16, 185, 129) # #10B981
    COLOR_GREEN_BG = RGBColor(236, 253, 245) # #ECFDF5
    COLOR_AMBER = RGBColor(217, 119, 6)       # #D97706
    COLOR_AMBER_BG = RGBColor(254, 243, 199)   # #FEF3C7
    COLOR_RED_ACCENT = RGBColor(225, 29, 72)   # #E11D48
    COLOR_RED_BG = RGBColor(255, 241, 242)    # #FFF1F2

    def add_header(slide, title_text, category_text="CX INTELLIGENCE & OPERATIONS PLATFORM"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_PRIMARY_BLUE

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN

        # Decorative subtle line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(226, 232, 240)
        line.line.color.rgb = RGBColor(226, 232, 240)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK_BG
    bg1.line.fill.background()

    # Decorative accent card on slide 1
    accent_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.1))
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    accent_box.line.color.rgb = RGBColor(51, 65, 85)

    # Tag pill
    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.7), Inches(3.2), Inches(0.4))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
    tag_box.line.fill.background()
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "BÁO CÁO NGHIỆP VỤ & ĐỊNH HƯỚNG"
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_WHITE
    p_tag.alignment = PP_ALIGN.CENTER

    # Main Title on Slide 1
    title_s1 = slide1.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(10.7), Inches(1.4))
    tf_s1 = title_s1.text_frame
    tf_s1.word_wrap = True
    p1_s1 = tf_s1.paragraphs[0]
    p1_s1.text = "CX Journey, Service & Hotspot Intelligence"
    p1_s1.font.size = Pt(32)
    p1_s1.font.bold = True
    p1_s1.font.color.rgb = COLOR_WHITE

    p2_s1 = tf_s1.add_paragraph()
    p2_s1.text = "Nền tảng Quản trị Trải nghiệm Cư dân & Điểm nóng Vận hành Bất động sản"
    p2_s1.font.size = Pt(18)
    p2_s1.font.color.rgb = RGBColor(148, 163, 184)

    # Subtitle / Scope
    desc_s1 = slide1.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10.7), Inches(1.0))
    tf_desc = desc_s1.text_frame
    tf_desc.word_wrap = True
    p_d1 = tf_desc.paragraphs[0]
    p_d1.text = "Báo cáo Khảo sát Nghiệp vụ Chuẩn hóa (Taxonomy 3.0.0) • Đề xuất Phạm vi Thử nghiệm Pilot (P0)\nTrình duyệt 7 Quyết định Nghiệp vụ Then chốt dành cho Ban Lãnh đạo"
    p_d1.font.size = Pt(14)
    p_d1.font.color.rgb = RGBColor(203, 213, 225)

    # Footer on Slide 1
    footer_s1 = slide1.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10.7), Inches(0.5))
    tf_f = footer_s1.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "Báo cáo viên: CX & Data Analyst Team  |  Dự án: Pilot CX Platform  |  Ngày: 17/08/2026"
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = RGBColor(100, 116, 139)

    # -------------------------------------------------------------
    # SLIDE 2: Bối cảnh & Vấn đề Cốt lõi (The Problem & The Shift)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "1. Bối cảnh & Mục tiêu Chuyển dịch Nghiệp vụ")

    # Left Card: Hiện trạng cũ
    card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_RED_BG
    card_left.line.color.rgb = RGBColor(254, 205, 211)

    tb_l = slide2.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "HIỆN TRẠNG & HẠN CHẾ CŨ"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_ACCENT

    items_left = [
        ("Chỉ đếm số lượng ticket/khiếu nại:", "Báo cáo chỉ cho biết 'có 500 ticket', không rõ khách bức xúc ở giai đoạn nào trong hành trình."),
        ("Dữ liệu khảo sát tĩnh trên Excel:", "Rời rạc, không liên kết trực tiếp với năng lực cung cấp dịch vụ thực tế của Ban Quản Lý (BQL)."),
        ("Xử lý sự vụ, bị động:", "Cùng một sự cố (như kẹt thang máy, ngập bãi xe) lặp lại nhiều lần tại 1 tòa nhà mà không có cảnh báo sớm."),
        ("Chưa tìm nguyên nhân gốc (RCA):", "Chỉ đóng ticket tạm thời khi cư dân phàn nàn, dẫn đến vấn đề liên tục tái diễn.")
    ]
    for title, desc in items_left:
        p_t = tf_l.add_paragraph()
        p_t.text = "❌ " + title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(8)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # Right Card: Định hướng Nền tảng CX Mới
    card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_GREEN_BG
    card_right.line.color.rgb = RGBColor(167, 243, 208)

    tb_r = slide2.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "MỤC TIÊU NỀN TẢNG CX MỚI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS_GREEN

    items_right = [
        ("Định vị đa chiều thông minh:", "Khách ở bước nào trong 36 bước hành trình → Sử dụng dịch vụ gì → Gặp sự cố cụ thể nào."),
        ("Gắn chặt với 10 Dịch vụ BQL:", "Xác định rõ ràng Bộ phận/Chuyên viên (Service Owner) chịu trách nhiệm giải quyết."),
        ("Phát hiện Điểm nóng (Hotspots):", "Tự động gom cụm theo không gian & thời gian (VD: Tòa S2 có 4 phản ánh thang máy/tuần)."),
        ("Xử lý nguyên nhân gốc (RCA):", "Liên kết bằng chứng thực tế, phân biệt triệu chứng với nguyên nhân để phòng ngừa triệt để.")
    ]
    for title, desc in items_right:
        p_t = tf_r.add_paragraph()
        p_t.text = "✅ " + title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(8)
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 3: Chu trình Vận hành Trải nghiệm (Closed-loop CX)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "2. Mô hình Vận hành Khép kín (Closed-Loop Operations)")

    steps = [
        ("1. LẮNG NGHE\n(Listen)", "Thu thập phản hồi từ App cư dân, Hotline, Lễ tân, Khảo sát", COLOR_PRIMARY_BLUE),
        ("2. THẤU HIỂU\n(Understand)", "Chuẩn hóa dữ liệu theo Taxonomy; AI gợi ý nhãn phân loại", COLOR_PRIMARY_BLUE),
        ("3. PHÁT HIỆN\n(Detect)", "Tự động gom cụm phát hiện Điểm nóng (Hotspot) theo Tòa/Khu", COLOR_AMBER),
        ("4. ƯU TIÊN\n(Prioritize)", "Phân cấp: Khẩn cấp (Immediate), Cần làm (Urgent), Kế hoạch", COLOR_AMBER),
        ("5. HÀNH ĐỘNG\n(Act)", "Điều phối đúng BQL / Kỹ thuật xử lý dứt điểm có SLA", COLOR_SUCCESS_GREEN),
        ("6. PHÒNG NGỪA\n(Learn)", "Phân tích nguyên nhân gốc rễ (RCA), tránh sự cố tái diễn", COLOR_SUCCESS_GREEN),
    ]

    for i, (title, desc, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(1.6 + row * 2.5)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = RGBColor(203, 213, 225)

        # Header bar in card
        hbar = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(0.6))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        tf_h = hbar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title.replace("\n", " - ")
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER

        # Body text
        tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(3.3), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # Bottom highlight bar
    bot_bar = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.6))
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    bot_bar.line.color.rgb = COLOR_BORDER_BLUE
    tf_b = bot_bar.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "⭐ NGUYÊN TẮC: Mọi chỉ số trên Dashboard đều có thể Bấm Xem Chi Tiết (Drill-down) đến từng phản hồi gốc có bằng chứng."
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_PRIMARY_BLUE
    p_b.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 4: Khung Nghiệp vụ Khảo sát Chuẩn hóa (Taxonomy 3.0.0)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "3. Khung Nghiệp vụ Khảo sát Chuẩn hóa (Taxonomy 3.0.0)")

    # Left Section: 2 Trục vòng đời
    card_life = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(4.5), Inches(5.2))
    card_life.fill.solid()
    card_life.fill.fore_color.rgb = COLOR_CARD_BG
    card_life.line.color.rgb = RGBColor(203, 213, 225)

    tb_lc = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.1), Inches(4.8))
    tf_lc = tb_lc.text_frame
    tf_lc.word_wrap = True
    p = tf_lc.paragraphs[0]
    p.text = "2 TRỤC VÒNG ĐỜI ĐỘC LẬP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    p1 = tf_lc.add_paragraph()
    p1.text = "1. Vòng đời Khách hàng (6 Giai đoạn / 36 Bước):"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    p1.space_before = Pt(6)

    stages = [
        "• Nhận thức (A1–A3): 3 bước tiếp cận ban đầu",
        "• Xem xét (C1–C6): 6 bước tư vấn & xem dự án",
        "• Giao dịch (TR01–TR06): 6 bước cọc & ký HĐMB",
        "• Nhận nhà (HO01–HO05): 5 bước nghiệm thu & nhận nhà",
        "• Cư trú (RES01–RES08): 8 bước đời sống cư dân",
        "• Vận hành (OPS01–OPS08): 8 bước quản lý của BQL"
    ]
    for st in stages:
        p_st = tf_lc.add_paragraph()
        p_st.text = st
        p_st.font.size = Pt(10)
        p_st.font.color.rgb = COLOR_TEXT_MUTED

    p2 = tf_lc.add_paragraph()
    p2.text = "2. Vòng đời Yêu cầu Dịch vụ (8 Bước):"
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(8)

    p_srv = tf_lc.add_paragraph()
    p_srv.text = "SRV-01 (Tìm tin) → SRV-02 (Gửi yêu cầu) → SRV-03 (Duyệt) → SRV-04 (Thanh toán) → SRV-05 (Phục vụ) → SRV-06 (Theo dõi) → SRV-07 (Hoàn tất) → SRV-08 (Đánh giá)."
    p_srv.font.size = Pt(10)
    p_srv.font.color.rgb = COLOR_TEXT_MUTED

    # Right Section: 10 Dịch vụ & 28 Sự cố
    card_svc = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(1.6), Inches(6.9), Inches(5.2))
    card_svc.fill.solid()
    card_svc.fill.fore_color.rgb = COLOR_CARD_BG
    card_svc.line.color.rgb = RGBColor(203, 213, 225)

    tb_sc = slide4.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(6.5), Inches(4.8))
    tf_sc = tb_sc.text_frame
    tf_sc.word_wrap = True
    p = tf_sc.paragraphs[0]
    p.text = "10 NHÓM DỊCH VỤ VẬN HÀNH & 28 SỰ CỐ CHUẨN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    services_list = [
        ("SV-01 Bán hàng & HĐ", "Tư vấn, đặt cọc, tiến độ HĐMB (3 Issues)"),
        ("SV-02 Thủ tục & Tài chính", "Phí quản lý, hóa đơn, công nợ (3 Issues)"),
        ("SV-03 Bàn giao nhà", "Lịch bàn giao, tồn tại nghiệm thu (3 Issues)"),
        ("SV-04 Hồ sơ cư dân & App", "Đăng ký cư dân, tính năng trên App (3 Issues)"),
        ("SV-05 Ra vào & Bãi xe", "Thẻ cư dân, cổng barie, chỗ đỗ xe (3 Issues)"),
        ("SV-06 Vệ sinh & Môi trường", "Rác thải, mùi hôi, cảnh quan cây xanh (3 Issues)"),
        ("SV-07 Kỹ thuật & Tiện ích", "Thang máy, điện nước chung, điều hòa (3 Issues)"),
        ("SV-08 An ninh & PCCC", "Bảo vệ, camera, báo cháy, an toàn (3 Issues)"),
        ("SV-09 Tiện ích trả phí", "Hồ bơi, sân tennis, gym, BBQ (3 Issues)"),
        ("SV-10 Dịch vụ khác", "Vấn đề phát sinh khác (Bắt buộc duyệt) (1 Issue)")
    ]

    for s_code, s_desc in services_list:
        p_s = tf_sc.add_paragraph()
        p_s.text = f"• {s_code}: {s_desc}"
        p_s.font.size = Pt(10.5)
        p_s.font.color.rgb = COLOR_TEXT_MAIN

    # -------------------------------------------------------------
    # SLIDE 5: Các Nguyên tắc Nghiệp vụ Bất biến (Kèm Ví dụ Thực tế)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "4. 4 Nguyên tắc Nghiệp vụ Bất biến (Kèm Ví dụ Thực tế)")

    principles = [
        ("1. HÀNH TRÌNH ≠ DỊCH VỤ (N:N)", 
         "Bước hành trình không đồng nhất với 1 dịch vụ duy nhất.",
         "VÍ DỤ THỰC TẾ:",
         "Khách ở bước 'RES-03 Ra vào & di chuyển' có thể phàn nàn về Dịch vụ Bãi xe (SV-05), Thang máy (SV-07), hoặc An ninh (SV-08).",
         COLOR_PRIMARY_BLUE),
        ("2. TRIỆU CHỨNG (ISSUE) ≠ NGUYÊN NHÂN (CAUSE)",
         "Sự cố khách nhìn thấy không phải là nguyên nhân kỹ thuật.",
         "VÍ DỤ THỰC TẾ:",
         "Khách báo 'Chờ thang máy quá lâu' (Triệu chứng/Issue); Nguyên nhân (Cause) có thể là do 'Hỏng biến tần' hoặc 'Giờ cao điểm học sinh đi học'.",
         COLOR_AMBER),
        ("3. GIẢ THUYẾT ≠ ROOT CAUSE XÁC NHẬN",
         "Phân định rõ giả định điều tra với kết luận có bằng chứng.",
         "VÍ DỤ THỰC TẾ:",
         "AI hoặc CSKH chỉ đưa ra Candidate Cause (Giả thuyết). Chỉ khi BQL/Kỹ thuật đến hiện trường kiểm tra có biên bản thì mới là Confirmed Root Cause.",
         COLOR_SUCCESS_GREEN),
        ("4. VAI TRÒ CỦA AI TRONG GIAI ĐOẠN PILOT",
         "AI đóng vai trò Trợ lý Gợi ý (Suggest-only).",
         "VÍ DỤ THỰC TẾ:",
         "Khi nạp file phản ánh, AI tự động gợi ý nhãn (Khách phàn nàn app cư dân). BQL/Reviewer bấm 'Duyệt' thì dữ liệu mới lên Dashboard chính thức.",
         COLOR_PRIMARY_BLUE)
    ]

    for i, (title, sub, ex_label, ex_desc, color) in enumerate(principles):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.6 + row * 2.6)

        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = RGBColor(203, 213, 225)

        # Header tag inside card
        tag = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(5.3), Inches(0.4))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf_t = tag.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

        # Body
        tb = slide5.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), Inches(5.2), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = sub
        p1.font.size = Pt(10.5)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN

        p2 = tf.add_paragraph()
        p2.text = f"📌 {ex_label} {ex_desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 6: Phân kỳ Triển khai: Pilot (P0) vs Mở rộng (P1)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "5. Phân kỳ Triển khai: Pilot (P0) vs Mở rộng (P1/P2)")

    # Left: P0 Pilot Scope
    card_p0 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_p0.fill.solid()
    card_p0.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    card_p0.line.color.rgb = COLOR_BORDER_BLUE

    tb_p0 = slide6.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_p0 = tb_p0.text_frame
    tf_p0.word_wrap = True
    p = tf_p0.paragraphs[0]
    p.text = "GIAI ĐOẠN 1: PILOT (P0) — TẬP TRUNG LÀM NGAY"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    p0_items = [
        ("Nạp & Chuẩn hóa Dữ liệu Lịch sử:", "Hỗ trợ import file CSV/Excel từ CSKH, Lễ tân; lọc trùng và che mờ thông tin cá nhân (PII Masking)."),
        ("Tách Phản ánh Đa ý định (Multi-intent):", "Cư dân gửi 1 câu chứa 2 việc (vừa chê app vừa kêu bãi xe) → Tách làm 2 item độc lập."),
        ("Không gian Duyệt nhãn AI (Feedback Workspace):", "AI gợi ý tự động, nhân viên duyệt nhanh 1-click để đưa dữ liệu vào báo cáo."),
        ("Thuật toán Điểm nóng (Hotspots P0):", "Tự động phát hiện cụm sự cố lặp lại theo Tòa nhà/Dịch vụ và gán người phụ trách."),
        ("4 Báo cáo Điều hành CX Cốt lõi:", "Cung cấp góc nhìn tổng quan cho Lãnh đạo & BQL Tòa nhà có tính năng Drill-down.")
    ]
    for t, d in p0_items:
        p_t = tf_p0.add_paragraph()
        p_t.text = "🔹 " + t
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(6)
        
        p_d = tf_p0.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # Right: P1 Expansion Scope
    card_p1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    card_p1.fill.solid()
    card_p1.fill.fore_color.rgb = COLOR_CARD_BG
    card_p1.line.color.rgb = RGBColor(203, 213, 225)

    tb_p1 = slide6.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_p1 = tb_p1.text_frame
    tf_p1.word_wrap = True
    p = tf_p1.paragraphs[0]
    p.text = "GIAI ĐOẠN 2: MỞ RỘNG (P1/P2) — SAU KHI PILOT XONG"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    p1_items = [
        ("Tích hợp API Thời gian thực:", "Kết nối trực tiếp với App Cư dân, Tổng đài Call Center, CRM, BMS/CMMS."),
        ("Quản trị Ticket & Giám sát SLA:", "Tạo ticket xử lý sự cố, theo dõi thời hạn cam kết SLA của từng nhà thầu vận hành."),
        ("Kích hoạt Cảnh báo Khẩn cấp (Hard Trigger):", "Tự động kích hoạt thông báo khẩn đối với sự cố an toàn/PCCC (SEV-1)."),
        ("Phân tích Xu hướng Nâng cao (MoM/YoY):", "So sánh biến động theo tháng và cùng kỳ năm trước khi có đủ 13 tháng dữ liệu thực."),
        ("Khảo sát Định kỳ Tự động (CSAT/CES/NPS):", "Tự động gửi khảo sát đo lường độ hài lòng sau mỗi lần cư dân được phục vụ.")
    ]
    for t, d in p1_items:
        p_t = tf_p1.add_paragraph()
        p_t.text = "🔸 " + t
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(6)
        
        p_d = tf_p1.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 7: Bức tranh 4 Dashboard Điều hành CX
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "6. Hệ thống 4 Dashboard Điều hành CX dành cho Lãnh đạo & BQL")

    dashboards = [
        ("1. CX OVERVIEW (Tổng quan)", 
         "Theo dõi 'nhịp tim' toàn diện của dự án", 
         ["• Tổng khối lượng phản ánh theo thời gian", "• Tỷ lệ hài lòng / tiêu cực / chưa rõ", "• Số lượng Điểm nóng đang hoạt động", "• Tỷ lệ hoàn thiện dữ liệu theo Tòa nhà"],
         COLOR_PRIMARY_BLUE),
        ("2. CUSTOMER JOURNEY (Hành trình)", 
         "Đo lường mức độ hài lòng theo 6 giai đoạn", 
         ["• Nhìn rõ khách đang 'tắc' ở bước nào nhất", "• So sánh mức độ tiêu cực giữa các giai đoạn", "• Điểm nghẽn từ Giao dịch đến Cư trú", "• Bấm vào bước để xem chi tiết phản ánh"],
         COLOR_PRIMARY_BLUE),
        ("3. SERVICE & PAIN POINTS (Dịch vụ)", 
         "Đánh giá chất lượng 10 Dịch vụ BQL", 
         ["• Bảng xếp hạng Dịch vụ & 28 Sự cố", "• Dịch vụ nào nhận nhiều phàn nàn nhất", "• Lọc linh hoạt theo Tòa nhà / Thời gian", "• So sánh hiệu quả vận hành giữa các tòa"],
         COLOR_AMBER),
        ("4. HOTSPOT ACTION QUEUE (Điểm nóng)", 
         "Bảng điều phối hành động dành cho BQL", 
         ["• Danh sách cụm sự cố lặp lại cần xử lý ngay", "• Phân loại: Khẩn cấp / Cần làm / Theo dõi", "• Gán người phụ trách (Service Owner)", "• Xem toàn bộ bằng chứng phản ánh liên quan"],
         COLOR_SUCCESS_GREEN),
    ]

    for i, (title, sub, bullets, color) in enumerate(dashboards):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.6 + row * 2.6)

        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = RGBColor(203, 213, 225)

        # Header tag
        tag = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(5.3), Inches(0.4))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf_t = tag.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

        # Content
        tb = slide7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), Inches(5.2), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_sub = tf.paragraphs[0]
        p_sub.text = sub
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_TEXT_MAIN

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = COLOR_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 8: 7 Quyết định Nghiệp vụ Then chốt cần Sếp Phê duyệt
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "7. Các Quyết định Nghiệp vụ Mở cần Sếp Chỉ đạo & Phê duyệt", "TRỌNG TÂM CUỘC HỌP")

    # Table of 7 Decisions
    rows = 8
    cols = 4
    table_shape = slide8.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    table = table_shape.table
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(4.5)
    table.columns[3].width = Inches(3.233)

    headers = ["STT", "Nội dung cần Sếp Phê duyệt", "Đề xuất Cụ thể của Team", "Tác động / Lý do cần chốt"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT

    decisions_data = [
        ("1", "Phạm vi Pilot (Scope)", "Chọn 1 Dự án cụ thể (VD: Tòa S2) & 3 Dịch vụ trọng điểm (Thang máy, Bãi xe, App).", "Khoanh vùng để kiểm soát chất lượng."),
        ("2", "Nguồn Dữ liệu Mẫu", "Trích xuất dữ liệu phản ánh/khảo sát thực tế 3-6 tháng gần nhất từ CSKH/Lễ tân.", "Có dữ liệu thật để kiểm thử mô hình."),
        ("3", "Phê duyệt Taxonomy 3.0.0", "Ban hành chính thức 10 Dịch vụ & Gán Service Owner cho từng Trưởng bộ phận.", "Xác định rõ người chịu trách nhiệm."),
        ("4", "Chính sách Bảo mật PII", "Che mờ SĐT, số căn khi AI xử lý; Chỉ Trưởng BQL/Admin được xem thông tin gốc.", "Đảm bảo tuân thủ bảo mật cư dân."),
        ("5", "Ngưỡng Điểm nóng (Hotspot)", "Cùng 1 lỗi xuất hiện ≥ 3 lần trong 7 ngày tại 1 Tòa nhà → Tạo cảnh báo Hotspot.", "Tránh báo động giả nhưng không lọt lỗi."),
        ("6", "Quy tắc Tách phản ánh", "Phản ánh chứa 2 vấn đề khác nhau sẽ tự động tách thành 2 việc riêng biệt.", "Phân loại chính xác cho từng bộ phận."),
        ("7", "Nhân sự tham gia Pilot", "Cử 1 CX Analyst, 1 Reviewer duyệt nhãn, và BQL Tòa Pilot tiếp nhận Hotspot.", "Đảm bảo có nhân sự vận hành thực tế.")
    ]

    for i, row in enumerate(decisions_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE if i % 2 == 0 else COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_TEXT_MAIN
            if j == 0:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 9: Kế hoạch Hành động Tiếp theo (Next Steps)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "8. Kế hoạch Hành động Tiếp theo (Next Steps & Roadmap)")

    phases = [
        ("TUẦN 1: THỐNG NHẤT", "Họp Workshop 1 buổi với các Trưởng bộ phận (BQL, CSKH, Kỹ thuật) để chốt 7 quyết định nghiệp vụ.", COLOR_PRIMARY_BLUE),
        ("TUẦN 2: CHUẨN BỊ", "Thu thập & làm sạch dữ liệu mẫu; Thiết lập hệ thống, tài khoản và phân quyền bảo mật PII.", COLOR_PRIMARY_BLUE),
        ("TUẦN 3-4: VẬN HÀNH PILOT", "Chạy thử nghiệm hệ thống; Đội ngũ duyệt nhãn và BQL tiếp nhận, xử lý Điểm nóng thực tế.", COLOR_AMBER),
        ("TUẦN 5: BÁO CÁO & MỞ RỘNG", "Đánh giá kết quả Pilot, độ chính xác của AI và trình Sếp kế hoạch kết nối API (Giai đoạn P1).", COLOR_SUCCESS_GREEN)
    ]

    for i, (title, desc, color) in enumerate(phases):
        x = Inches(0.8 + i * 2.95)
        y = Inches(1.6)
        
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = RGBColor(203, 213, 225)

        hbar = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(0.6))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        tf_h = hbar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER

        tb = slide9.shapes.add_textbox(x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(2.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # Bottom Proposal Box
    prop_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.7))
    prop_box.fill.solid()
    prop_box.fill.fore_color.rgb = COLOR_GREEN_BG
    prop_box.line.color.rgb = RGBColor(167, 243, 208)

    tb_prop = slide9.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(11.1), Inches(1.5))
    tf_prop = tb_prop.text_frame
    tf_prop.word_wrap = True
    p_pr = tf_prop.paragraphs[0]
    p_pr.text = "🎯 KIẾN NGHỊ SẾP PHÊ DUYỆT TRONG BUỔI HÔM NAY:"
    p_pr.font.size = Pt(12)
    p_pr.font.bold = True
    p_pr.font.color.rgb = COLOR_SUCCESS_GREEN

    p_a1 = tf_prop.add_paragraph()
    p_a1.text = "1. Phê duyệt chủ trương triển khai thử nghiệm Giai đoạn Pilot (P0) theo khung Taxonomy 3.0.0."
    p_a1.font.size = Pt(11)
    p_a1.font.bold = True
    p_a1.font.color.rgb = COLOR_TEXT_MAIN
    p_a1.space_before = Pt(4)

    p_a2 = tf_prop.add_paragraph()
    p_a2.text = "2. Đồng ý cho phép tổ chức 01 buổi Workshop ngắn với các Trưởng bộ phận liên quan để chốt chi tiết 7 Quyết định mở."
    p_a2.font.size = Pt(11)
    p_a2.font.bold = True
    p_a2.font.color.rgb = COLOR_TEXT_MAIN
    p_a2.space_before = Pt(2)

    # Save presentation
    output_path = "/Users/thangnguyen/Documents/analyst-data-workspace/Bao_Cao_Nghiep_Vu_CX_Platform.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
